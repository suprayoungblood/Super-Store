"""PowerPoint rendering of the slide specs (the 'how').

Each :mod:`superstore.slides` spec is rendered by a dedicated handler, selected
through a type-to-renderer mapping so there is no long ``if/elif`` chain.
``python-pptx`` is imported lazily so the rest of the package works even when
that optional dependency is not installed.
"""

from __future__ import annotations

from pathlib import Path

from .analysis import AnalysisResult
from .config import AppConfig
from .slides import BulletsSlide, ImageSlide, TitleSlide, build_slides

# Indices into the default template's slide layouts.
_TITLE_LAYOUT = 0
_CONTENT_LAYOUT = 1
_TITLE_ONLY_LAYOUT = 5

# Shared slide layout metrics, in inches. Kept as raw floats so the ``Inches``
# converter is applied only after the lazy ``python-pptx`` import in each helper.
_MARGIN_INCHES = 0.5
_IMAGE_TOP_INCHES = 1.5
_CAPTION_HEIGHT_INCHES = 0.8


def build_deck(
    result: AnalysisResult, versions: dict[str, str], config: AppConfig
) -> Path:
    """Render the 7-slide report to a ``.pptx`` file.

    Args:
        result: Computed metrics used to populate the slides.
        versions: Installed library versions for the installation slide.
        config: Provides presenter metadata and output paths.

    Returns:
        Path to the saved PowerPoint file.

    Raises:
        RuntimeError: If ``python-pptx`` is not installed.
        FileNotFoundError: If the chart image has not been rendered yet.
    """
    presentation_factory = _import_presentation()
    _require_chart(config.paths.chart_file)
    deck = presentation_factory()
    for spec in build_slides(result, versions, config):
        _RENDERERS[type(spec)](deck, spec)
    output = config.paths.deck_file
    output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(str(output))
    return output


def _import_presentation():
    """Lazily import ``python-pptx``'s ``Presentation`` with a clear error."""
    try:
        from pptx import Presentation
    except ImportError as error:
        raise RuntimeError(
            "python-pptx is required to build the deck. Install it with "
            "'pip install python-pptx'."
        ) from error
    return Presentation


def _require_chart(path: Path) -> None:
    """Raise if the chart image to embed is missing."""
    if not path.is_file():
        raise FileNotFoundError(
            f"Chart image '{path}' not found. Render the chart before building "
            f"the deck."
        )


def _add_slide(deck, layout_index: int):
    """Append a slide using the given layout and return it."""
    return deck.slides.add_slide(deck.slide_layouts[layout_index])


def _render_title(deck, spec: TitleSlide) -> None:
    """Render the title slide (heading + multi-line subtitle)."""
    slide = _add_slide(deck, _TITLE_LAYOUT)
    slide.shapes.title.text = spec.title
    slide.placeholders[1].text = spec.subtitle


def _render_bullets(deck, spec: BulletsSlide) -> None:
    """Render a titled slide whose body is a bulleted list."""
    slide = _add_slide(deck, _CONTENT_LAYOUT)
    slide.shapes.title.text = spec.title
    body = slide.placeholders[1].text_frame
    body.clear()
    for index, bullet in enumerate(spec.bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet


def _render_image(deck, spec: ImageSlide) -> None:
    """Render a titled slide that embeds a centered image with a caption."""
    slide = _add_slide(deck, _TITLE_ONLY_LAYOUT)
    slide.shapes.title.text = spec.title
    _place_centered_image(deck, slide, spec.image_path)
    _add_caption(deck, slide, spec.caption)


def _place_centered_image(deck, slide, image_path: Path) -> None:
    """Scale an image to fit below the title and center it horizontally."""
    from pptx.util import Inches

    margin = Inches(_MARGIN_INCHES)
    top = Inches(_IMAGE_TOP_INCHES)
    caption_height = Inches(_CAPTION_HEIGHT_INCHES)
    available_width = deck.slide_width - 2 * margin
    available_height = deck.slide_height - top - caption_height - margin

    picture = slide.shapes.add_picture(str(image_path), margin, top)
    scale = min(available_width / picture.width, available_height / picture.height)
    picture.width = int(picture.width * scale)
    picture.height = int(picture.height * scale)
    picture.left = int((deck.slide_width - picture.width) / 2)
    picture.top = top


def _add_caption(deck, slide, caption: str) -> None:
    """Add an italic caption textbox near the bottom of the slide."""
    from pptx.util import Inches, Pt

    margin = Inches(_MARGIN_INCHES)
    caption_height = Inches(_CAPTION_HEIGHT_INCHES)
    box = slide.shapes.add_textbox(
        margin,
        deck.slide_height - caption_height - margin,
        deck.slide_width - 2 * margin,
        caption_height,
    )
    box.text_frame.word_wrap = True
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = caption
    paragraph.font.size = Pt(14)
    paragraph.font.italic = True


# Type-to-renderer dispatch table (defined after the handlers exist).
_RENDERERS = {
    TitleSlide: _render_title,
    BulletsSlide: _render_bullets,
    ImageSlide: _render_image,
}
